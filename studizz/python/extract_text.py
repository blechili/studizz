#!/usr/bin/env python3
"""
STUDIZ — Text Extraction Script
---------------------------------
Usage:
    python3 extract_text.py /absolute/path/to/file.upload <original_extension>

Called exclusively by PHP's proc_open with a validated absolute file path
(always stored on disk with a fixed, inert ".upload" extension — see
api_gateway.php) plus the client's ORIGINAL file extension, passed as a
plain string used only to pick which parser to try. Writes extracted plain
text to stdout. Writes diagnostic messages to stderr. Exits with code 0 on
success (including "no usable text found", which intentionally still exits
0 with empty output — see main()), non-zero only on a genuine extraction
error (corrupted file, missing dependency, etc).

Supported input types:
  - PDF            (digital text vectors and scanned/image-based pages)
  - JPEG / PNG     (OCR via Tesseract)
  - DOCX           (Word documents, including text inside tables)
  - PPTX           (PowerPoint 2007+, including notes and table text)
  - PPT            (legacy PowerPoint 97-2003 binary format — BEST EFFORT.
                     No parser library is used; this scans the raw binary
                     for UTF-16LE text runs with noise filtering. Captures
                     most real slide text but isn't a real parser: slide
                     order isn't preserved and a little non-content noise
                     can slip through. See extract_from_legacy_ppt().)
  - XLSX           (Excel 2007+ workbooks, via openpyxl — every sheet's
                     cells rendered as pipe-separated rows)
  - XLS            (legacy Excel 97-2003 binary format, via xlrd — a real
                     parser, not a heuristic; xlrd still supports reading
                     this older format even though it dropped xlsx support)
  - TXT / MD / CSV / LOG  (read directly as text)

Any other extension is treated as having no extractable text rather than
an error — the caller (api_gateway.php) turns that into a friendly
"not enough readable text" rejection instead of a crash.

Dependencies (install on server):
    pip install pymupdf pytesseract Pillow python-docx python-pptx openpyxl xlrd
    apt-get install tesseract-ocr  (for OCR on images and scanned PDFs)
"""

import sys
import os
import re
import datetime as _dt
import traceback

# Force UTF-8 output so Windows cp1252 doesn't choke on special characters
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

def log_err(msg: str) -> None:
    """Write a timestamped diagnostic line to stderr."""
    print(f"[extract_text] {msg}", file=sys.stderr, flush=True)


def extract_from_pdf(path: str) -> str:
    """
    Extract text from a PDF file using PyMuPDF (fitz).
    For pages that contain no selectable text (scanned/image pages),
    fall back to Tesseract OCR via pytesseract.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        log_err("PyMuPDF (fitz) is not installed. Run: pip install pymupdf")
        sys.exit(2)

    doc = fitz.open(path)
    all_text: list[str] = []
    ocr_pages: list[int] = []

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            all_text.append(text)
        else:
            # No selectable text — mark for OCR
            ocr_pages.append((page_num, page))

    if ocr_pages:
        log_err(f"{len(ocr_pages)} page(s) require OCR.")
        try:
            import pytesseract
            from PIL import Image
            import io
        except ImportError:
            log_err("pytesseract or Pillow not installed. Scanned pages will be skipped.")
            ocr_pages = []

        for page_num, page in ocr_pages:
            try:
                # Render page to a high-resolution pixel map
                mat       = fitz.Matrix(2.0, 2.0)
                pix       = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
                img_bytes = pix.tobytes("png")
                image     = Image.open(io.BytesIO(img_bytes))
                ocr_text  = pytesseract.image_to_string(image, lang="eng")
                if ocr_text.strip():
                    all_text.append(f"[Page {page_num} — OCR]\n{ocr_text.strip()}")
                else:
                    log_err(f"Page {page_num}: OCR returned no text.")
            except Exception as e:
                log_err(f"Page {page_num}: OCR failed, skipping. {e}")
                continue
    doc.close()
    return "\n\n".join(all_text)


def extract_from_image(path: str) -> str:
    """
    Extract text from an image file (JPG / PNG) using Tesseract OCR.
    """
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        log_err("pytesseract or Pillow not installed. Run: pip install pytesseract Pillow")
        sys.exit(2)

    image = Image.open(path)

    # Convert to RGB if necessary (handles RGBA / palette modes)
    if image.mode not in ("RGB", "L"):
        image = image.convert("RGB")

    # Upscale small images to improve OCR accuracy
    width, height = image.size
    if width < 1000 or height < 1000:
        scale  = max(1000 / width, 1000 / height)
        image  = image.resize(
            (int(width * scale), int(height * scale)),
            resample=Image.LANCZOS
        )

    text = pytesseract.image_to_string(image, lang="eng", config="--psm 6")
    return text.strip()


def extract_from_docx(path: str) -> str:
    """
    Extract text from a Word .docx file, including paragraph text and
    any text found inside tables.
    """
    try:
        import docx
    except ImportError:
        log_err("python-docx is not installed. Run: pip install python-docx")
        sys.exit(2)

    document = docx.Document(path)
    parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]

    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text.strip())

    return "\n\n".join(parts)


def extract_from_pptx(path: str) -> str:
    """
    Extract text from a modern PowerPoint .pptx file: every shape's text
    frame, table cells, and speaker notes, slide by slide.
    """
    try:
        from pptx import Presentation
    except ImportError:
        log_err("python-pptx is not installed. Run: pip install python-pptx")
        sys.exit(2)

    prs = Presentation(path)
    slides_text: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text.strip())
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker notes] {notes}")
        if parts:
            slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(parts))

    return "\n\n".join(slides_text)


def _cell_to_str(value) -> str:
    """
    Render a single spreadsheet cell value as readable text. Whole-number
    floats print without a trailing ".0", and dates print as plain
    YYYY-MM-DD (or with a time component if one is set) instead of a
    raw Excel serial number — both are common footguns that would
    otherwise confuse "key figures" in the cell with junk.
    """
    if value is None:
        return ""
    if isinstance(value, _dt.datetime):
        if value.time() == _dt.time():
            return value.strftime("%Y-%m-%d")
        return value.strftime("%Y-%m-%d %H:%M")
    if isinstance(value, _dt.date):
        return value.strftime("%Y-%m-%d")
    if isinstance(value, float) and value == int(value):
        return str(int(value))
    return str(value).strip()


def extract_from_xlsx(path: str) -> str:
    """
    Extract text from a modern Excel .xlsx workbook: every sheet's cells
    rendered as simple pipe-separated rows so the AI can read rows/columns
    and spot headers and key figures. Computed values are used in place
    of formulas (data_only=True). Capped at MAX_ROWS_PER_SHEET rows per
    sheet to bound memory/time on huge spreadsheets.
    """
    try:
        import openpyxl
    except ImportError:
        log_err("openpyxl is not installed. Run: pip install openpyxl")
        sys.exit(2)

    MAX_ROWS_PER_SHEET = 2000
    sheets_text: list[str] = []

    # openpyxl rejects a string path outright unless it ends in .xlsx/.xlsm/
    # .xltx/.xltm — which our on-disk file never does (always a fixed,
    # inert ".upload" extension, by design; see api_gateway.php). Passing
    # a file object instead skips that extension check entirely (openpyxl
    # only applies it to string paths) and goes straight to reading the
    # zip contents, which don't care about the filename. Stays open for
    # the whole read since read_only mode streams rows lazily from it.
    with open(path, "rb") as f:
        wb = openpyxl.load_workbook(f, data_only=True, read_only=True)
        for ws in wb.worksheets:
            rows_out: list[str] = []
            truncated = False
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= MAX_ROWS_PER_SHEET:
                    truncated = True
                    break
                cells = [_cell_to_str(c) for c in row]
                if not any(cells):
                    continue
                rows_out.append(" | ".join(cells))
            if not rows_out:
                continue
            if truncated:
                rows_out.append(f"... ({MAX_ROWS_PER_SHEET}-row limit reached, remaining rows omitted)")
            sheets_text.append(f"[Sheet: {ws.title}]\n" + "\n".join(rows_out))
        wb.close()

    return "\n\n".join(sheets_text)


def extract_from_xls(path: str) -> str:
    """
    Extract text from a legacy Excel .xls workbook (Excel 97-2003 binary
    format) using xlrd, which still supports reading this older format
    (it dropped .xlsx support in v2.0 but kept .xls). Same pipe-separated
    row rendering as extract_from_xlsx — a real parser, not a heuristic.
    """
    try:
        import xlrd
    except ImportError:
        log_err("xlrd is not installed. Run: pip install xlrd")
        sys.exit(2)

    MAX_ROWS_PER_SHEET = 2000

    book = xlrd.open_workbook(path)
    sheets_text: list[str] = []

    for sheet in book.sheets():
        rows_out: list[str] = []
        truncated = sheet.nrows > MAX_ROWS_PER_SHEET
        for r in range(min(sheet.nrows, MAX_ROWS_PER_SHEET)):
            cells = []
            for c in range(sheet.ncols):
                cell = sheet.cell(r, c)
                if cell.ctype == xlrd.XL_CELL_DATE:
                    try:
                        value = xlrd.xldate_as_datetime(cell.value, book.datemode)
                    except Exception:
                        value = cell.value
                else:
                    value = cell.value
                cells.append(_cell_to_str(value))
            if not any(cells):
                continue
            rows_out.append(" | ".join(cells))
        if not rows_out:
            continue
        if truncated:
            rows_out.append(f"... ({MAX_ROWS_PER_SHEET}-row limit reached, remaining rows omitted)")
        sheets_text.append(f"[Sheet: {sheet.name}]\n" + "\n".join(rows_out))

    return "\n\n".join(sheets_text)


_PPT_TEXT_RUN_RE      = re.compile(rb'(?:[\x09\x0A\x0D\x20-\x7E]\x00){4,}')
_PPT_NOISE_EXACT       = {
    "root entry", "current user", "summaryinformation", "powerpoint document",
    "documentsummaryinformation", "kso_wm_beautify_flag",
}
_PPT_FONT_NAMES = {
    "tahoma", "arial", "wingdings", "simsun", "calibri", "cambria", "verdana",
    "times new roman", "courier new", "georgia", "segoe ui", "+mn-ea", "+mj-lt", "+mn-lt",
}
_PPT_SHAPE_NAME_RE  = re.compile(r'^[A-Za-z][A-Za-z\s]*\s\d+$')   # "Title 3073", "Rectangles 46082"
_PPT_LAYOUT_NAME_RE = re.compile(r'^\d+_\S.*$')                  # "1_Curtain Call"
_PPT_INTERNAL_RE    = re.compile(r'^_{2,}')                      # "___PPT12", "__WPP10"


def _ppt_looks_like_noise(line: str) -> bool:
    s = line.strip()
    if not s:
        return True
    low = s.lower()
    if low in _PPT_NOISE_EXACT or low in _PPT_FONT_NAMES:
        return True
    if _PPT_INTERNAL_RE.match(s):
        return True
    if _PPT_SHAPE_NAME_RE.match(s):
        return True
    if _PPT_LAYOUT_NAME_RE.match(s):
        return True
    return False


def extract_from_legacy_ppt(path: str) -> str:
    """
    BEST-EFFORT extraction for legacy .ppt (PowerPoint 97-2003, an
    OLE2/CFBF binary format). No dedicated parser library is used — this
    scans the raw bytes for contiguous runs of UTF-16LE-decodable text,
    which is how the format stores most slide text, then filters out
    recognisable non-content noise (shape names like "Title 3073", font
    names, internal stream names). Lower fidelity than a real parser:
    slide order isn't preserved and a little noise can still slip
    through — but it needs no extra dependencies and works on any OS.
    """
    with open(path, "rb") as f:
        raw = f.read()

    runs: list[str] = []
    seen: set[str] = set()
    for match in _PPT_TEXT_RUN_RE.finditer(raw):
        chunk = match.group(0).decode("utf-16-le", errors="ignore").strip()
        if len(chunk) < 3:
            continue
        key = chunk.lower()
        if key in seen:
            continue
        seen.add(key)
        if _ppt_looks_like_noise(chunk):
            continue
        runs.append(chunk)

    return "\n".join(runs)


def extract_from_text(path: str) -> str:
    """
    Read a plain-text-like file (.txt, .md, .csv, .log) directly.
    Falls back to lossy decoding rather than failing on odd encodings.
    """
    with open(path, "rb") as f:
        raw = f.read()
    try:
        return raw.decode("utf-8-sig")  # -sig strips a leading BOM if present
    except UnicodeDecodeError:
        return raw.decode("utf-8", errors="replace")


def detect_type(ext: str) -> str:
    """
    Classify the file purely from the client's original extension
    (the on-disk path is always a fixed ".upload" name and tells us
    nothing). Returns 'pdf', 'image', 'docx', 'pptx', 'ppt_legacy',
    'xlsx', 'xls_legacy', 'text', or 'unsupported'.
    """
    ext = ext.lower().lstrip(".")

    if ext == "pdf":
        return "pdf"
    if ext in ("jpg", "jpeg", "png"):
        return "image"
    if ext == "docx":
        return "docx"
    if ext == "pptx":
        return "pptx"
    if ext == "ppt":
        return "ppt_legacy"
    if ext == "xlsx":
        return "xlsx"
    if ext == "xls":
        return "xls_legacy"
    if ext in ("txt", "md", "csv", "log"):
        return "text"
    return "unsupported"


def main() -> None:
    if len(sys.argv) != 3:
        log_err("Usage: extract_text.py <absolute_file_path> <original_extension>")
        sys.exit(1)

    file_path = sys.argv[1]
    orig_ext  = sys.argv[2]

    # Security: reject any path that isn't absolute
    if not os.path.isabs(file_path):
        log_err(f"Rejected relative path: {file_path}")
        sys.exit(1)

    # Security: ensure the file exists
    if not os.path.isfile(file_path):
        log_err(f"File not found: {file_path}")
        sys.exit(1)

    file_type = detect_type(orig_ext)
    log_err(f"Processing as {file_type.upper()} (original ext: .{orig_ext})")

    if file_type == "unsupported":
        # Not an error — this file type simply has no text we know how to
        # pull out of it. Exit 0 with empty output so the caller's normal
        # "not enough readable text" rejection handles it with one
        # consistent, friendly message rather than a hard crash.
        log_err("No extractor available for this file type — treating as empty.")
        sys.stdout.write("")
        sys.stdout.flush()
        sys.exit(0)

    try:
        if file_type == "pdf":
            text = extract_from_pdf(file_path)
        elif file_type == "image":
            text = extract_from_image(file_path)
        elif file_type == "docx":
            text = extract_from_docx(file_path)
        elif file_type == "pptx":
            text = extract_from_pptx(file_path)
        elif file_type == "ppt_legacy":
            text = extract_from_legacy_ppt(file_path)
        elif file_type == "xlsx":
            text = extract_from_xlsx(file_path)
        elif file_type == "xls_legacy":
            text = extract_from_xls(file_path)
        else:
            text = extract_from_text(file_path)
    except Exception as exc:  # noqa: BLE001
        log_err(f"Extraction failed: {exc}")
        log_err(traceback.format_exc())
        sys.exit(1)

    # Write extracted text to stdout for PHP to read. Deliberately no
    # length check here — even an empty/short result exits 0 and lets
    # PHP's single "not enough readable text" threshold be the one
    # source of truth for that message.
    sys.stdout.write(text or "")
    sys.stdout.flush()
    sys.exit(0)


if __name__ == "__main__":
    main()
